# lot of stuff is copied from https://github.com/QwenLM/Qwen-Agent/blob/8bfb20e6d10cd4f1f0250eeda995692cbb76dd47/qwen_agent/tools/simple_doc_parser.py#L31
# all credit goes to Qwen team for writing so much clean code which works reasonable for my use case
# I just adapted as per my needs ; not worth re-doing all the parsing stuff myself :)
# however, a lot of url types don't work currently; we need to support them as well

import logging
import os
import re
import urllib
from collections import Counter
from pathlib import Path
from typing import List

import pandas as pd
import pdfplumber
import requests
from bs4 import BeautifulSoup
from docx import Document
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTChar, LTRect, LTTextContainer
from pptx import Presentation
from tabulate import tabulate

from .utils import hash_sha256, make_request

logger = logging.getLogger(__name__)


def clean_paragraph(text):
    # remove cid
    text = re.sub(r"\(cid:\d+\)", "", text)

    # remove hexadecimal
    text = re.sub(r"[0-9A-Fa-f]{21,}", "", text)

    # remove continuous_placeholders
    text = re.sub(r"[.\- —。_*]{7,}", "\t", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text


PARAGRAPH_SPLIT_SYMBOL = "\n"


def parse_word(docx_path: str):
    doc = Document(docx_path)

    content = []
    for para in doc.paragraphs:
        content.append({"text": para.text})
    for table in doc.tables:
        tbl = []
        for row in table.rows:
            tbl.append("|" + "|".join([cell.text for cell in row.cells]) + "|")
        tbl = "\n".join(tbl)
        content.append({"table": tbl})

    # Due to the pages in Word are not fixed, the entire document is returned as one page
    return [{"page_num": 1, "content": content}]


def parse_ppt(path: str):
    ppt = Presentation(path)
    doc = []
    for slide_number, slide in enumerate(ppt.slides):
        page = {"page_num": slide_number + 1, "content": []}

        for shape in slide.shapes:
            if not shape.has_text_frame and not shape.has_table:
                pass

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    paragraph_text = clean_paragraph(paragraph_text)
                    if paragraph_text.strip():
                        page["content"].append({"text": paragraph_text})

            if shape.has_table:
                tbl = []
                for row_number, row in enumerate(shape.table.rows):
                    tbl.append("|" + "|".join([cell.text for cell in row.cells]) + "|")
                tbl = "\n".join(tbl)
                page["content"].append({"table": tbl})
        doc.append(page)
    return doc


def parse_txt(path: str):
    with open(path, "r", encoding="utf-8") as file:
        text = file.read()
    paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
    content = []
    for p in paras:
        content.append({"text": p})

    # Due to the pages in txt are not fixed, the entire document is returned as one page
    return [{"page_num": 1, "content": content}]


def df_to_md(df) -> str:

    def replace_long_dashes(text):
        if text.replace("-", "").replace(":", "").strip():
            return text
        pattern = r"-{6,}"
        replaced_text = re.sub(pattern, "-----", text)
        return replaced_text

    df = df.dropna(how="all")
    df = df.dropna(axis=1, how="all")
    df = df.fillna("")
    md_table = tabulate(df, headers="keys", tablefmt="pipe", showindex=False)

    md_table = "\n".join(
        [
            "|".join(
                replace_long_dashes(" " + cell.strip() + " " if cell else "")
                for cell in row.split("|")
            )
            for row in md_table.split("\n")
        ]
    )
    return md_table


def parse_excel(file_path: str) -> List[dict]:
    excel_file = pd.ExcelFile(file_path)
    md_tables = []
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        md_table = df_to_md(df)
        md_tables.append(f"### Sheet: {sheet_name}\n{md_table}")

    return [
        {"page_num": i + 1, "content": [{"table": md_tables[i]}]}
        for i in range(len(md_tables))
    ]


def parse_csv(file_path: str) -> List[dict]:
    md_tables = []
    try:
        df = pd.read_csv(file_path, encoding_errors="replace", on_bad_lines="skip")
    except:
        return parse_excel(file_path)
    md_table = df_to_md(df)
    md_tables.append(md_table)  # There is only one table available

    return [
        {"page_num": i + 1, "content": [{"table": md_tables[i]}]}
        for i in range(len(md_tables))
    ]


def parse_tsv(file_path: str) -> List[dict]:
    md_tables = []
    try:
        df = pd.read_csv(
            file_path, sep="\t", encoding_errors="replace", on_bad_lines="skip"
        )
    except:
        return parse_excel(file_path)
    md_table = df_to_md(df)
    md_tables.append(md_table)  # There is only one table available

    return [
        {"page_num": i + 1, "content": [{"table": md_tables[i]}]}
        for i in range(len(md_tables))
    ]


def parse_html_bs(path: str):
    def pre_process_html(s):
        # replace multiple newlines
        s = re.sub("\n+", "\n", s)
        # replace special string
        s = s.replace("Add to Qwen's Reading List", "")
        return s

    bs_kwargs = {"features": "lxml"}
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, **bs_kwargs)

    text = soup.get_text()

    if soup.title:
        title = str(soup.title.string)
    else:
        title = ""

    text = pre_process_html(text)
    paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
    content = []
    for p in paras:
        p = clean_paragraph(p)
        if p.strip():
            content.append({"text": p})

    # The entire document is returned as one page
    return [{"page_num": 1, "content": content, "title": title}]


def parse_pdf(pdf_path: str) -> List[dict]:
    # Todo: header and footer

    doc = []
    pdf = pdfplumber.open(pdf_path)
    for i, page_layout in enumerate(extract_pages(pdf_path)):
        page = {"page_num": page_layout.pageid, "content": []}

        elements = []
        for element in page_layout:
            elements.append(element)

        # Init params for table
        table_num = 0
        tables = []

        for element in elements:
            if isinstance(element, LTRect):
                if not tables:
                    tables = extract_tables(pdf, i)
                if table_num < len(tables):
                    table_string = table_converter(tables[table_num])
                    table_num += 1
                    if table_string:
                        page["content"].append({"table": table_string, "obj": element})
            elif isinstance(element, LTTextContainer):
                # Delete line breaks in the same paragraph
                text = element.get_text()
                # Todo: Further analysis using font
                font = get_font(element)
                if text.strip():
                    new_content_item = {"text": text, "obj": element}
                    if font:
                        new_content_item["font-size"] = round(font[1])
                        # new_content_item['font-name'] = font[0]
                    page["content"].append(new_content_item)
            else:
                # Other types such as Image or Figure are not supported for now
                pass

        # merge elements
        page["content"] = postprocess_page_content(page["content"])
        doc.append(page)

    return doc


def postprocess_page_content(page_content: list) -> list:
    # rm repetitive identification for table and text
    # Some documents may repeatedly recognize LTRect and LTTextContainer
    table_obj = [p["obj"] for p in page_content if "table" in p]
    tmp = []
    for p in page_content:
        repetitive = False
        if "text" in p:
            for t in table_obj:
                if (
                    t.bbox[0] <= p["obj"].bbox[0]
                    and p["obj"].bbox[1] <= t.bbox[1]
                    and t.bbox[2] <= p["obj"].bbox[2]
                    and p["obj"].bbox[3] <= t.bbox[3]
                ):
                    repetitive = True
                    break

        if not repetitive:
            tmp.append(p)
    page_content = tmp

    # merge paragraphs that have been separated by mistake
    new_page_content = []
    for p in page_content:
        if (
            new_page_content
            and "text" in new_page_content[-1]
            and "text" in p
            and abs(p.get("font-size", 12) - new_page_content[-1].get("font-size", 12))
            < 2
            and p["obj"].height < p.get("font-size", 12) + 1
        ):
            # Merge those lines belonging to a paragraph
            _p = p["text"]
            new_page_content[-1]["text"] += f" {_p}"
            # new_page_content[-1]['font-name'] = p.get('font-name', '')
            new_page_content[-1]["font-size"] = p.get("font-size", 12)
        else:
            p.pop("obj")
            new_page_content.append(p)
    for i in range(len(new_page_content)):
        if "text" in new_page_content[i]:
            new_page_content[i]["text"] = clean_paragraph(new_page_content[i]["text"])
    return new_page_content


def get_font(element):
    fonts_list = []
    for text_line in element:
        if isinstance(text_line, LTTextContainer):
            for character in text_line:
                if isinstance(character, LTChar):
                    fonts_list.append((character.fontname, character.size))

    fonts_list = list(set(fonts_list))
    if fonts_list:
        counter = Counter(fonts_list)
        most_common_fonts = counter.most_common(1)[0][0]
        return most_common_fonts
    else:
        return []


def extract_tables(pdf, page_num):
    table_page = pdf.pages[page_num]
    tables = table_page.extract_tables()
    return tables


def table_converter(table):
    table_string = ""
    for row_num in range(len(table)):
        row = table[row_num]
        cleaned_row = [
            (
                item.replace("\n", " ")
                if item is not None and "\n" in item
                else "None" if item is None else item
            )
            for item in row
        ]
        table_string += "|" + "|".join(cleaned_row) + "|" + "\n"
    table_string = table_string[:-1]
    return table_string


PARSER_SUPPORTED_FILE_TYPES = [
    "pdf",
    "docx",
    "pptx",
    "txt",
    "html",
    "csv",
    "tsv",
    "xlsx",
    "xls",
]


def is_http_url(url):
    return url.startswith("https://") or url.startswith("http://")


def contains_html_tags(text: str) -> bool:
    pattern = r"<(p|span|div|li|html|script)[^>]*?"
    return bool(re.search(pattern, text))


def get_basename_from_url(path_or_url: str) -> str:
    # "/mnt/a/b/c" -> "c"
    # "https://github.com/here?k=v" -> "here"
    # "https://github.com/" -> ""
    basename = urllib.parse.urlparse(path_or_url).path
    basename = os.path.basename(basename)
    basename = urllib.parse.unquote(basename)
    basename = basename.strip()
    # "https://github.com/" -> "" -> "github.com"
    if not basename:
        basename = [x.strip() for x in path_or_url.split("/") if x.strip()][-1]
    return basename


def get_content_type_by_head_request(path: str) -> str:
    try:
        response = requests.head(path, timeout=5)
        content_type = response.headers.get("Content-Type", "")
        return content_type
    except requests.RequestException:
        return "unk"


def get_file_type(path: str):
    f_type = get_basename_from_url(path).split(".")[-1].lower()
    if f_type in ["pdf", "docx", "pptx", "csv", "tsv", "xlsx", "xls"]:
        return f_type

    if is_http_url(path):
        # The HTTP header information for the response is obtained by making a HEAD request to the target URL,
        # where the Content-type field usually indicates the Type of Content to be returned
        content_type = get_content_type_by_head_request(path)
        if "application/pdf" in content_type:
            return "pdf"
        elif "application/msword" in content_type:
            return "docx"
        # Assuming that the URL is HTML by default,
        # because the file downloaded by the request may contain html tags
        return "html"
    else:
        # Determine by reading local HTML file
        try:
            with open(path, "r", encoding="utf-8") as file:
                content = file.read()
        except Exception:
            return "unk"
        if contains_html_tags(content):
            return "html"
        else:
            return "txt"


def get_plain_doc(doc: list):
    paras = []
    for page in doc:
        for para in page["content"]:
            for k, v in para.items():
                if k in ["text", "table", "image"]:
                    paras.append(v)
    return PARAGRAPH_SPLIT_SYMBOL.join(paras)


# TODO: Currently, these links don't work well - they need browser based stuff with playwright
# url = "https://www.linkedin.com/posts/charlesmartin14_knuth-is-now-vibe-mathing-shock-shock-activity-7434661528378859520-ziOR?utm_source=share&utm_medium=member_desktop&rcm=ACoAACm304YBKyRBusHCwziwrqJebOhHonRd5bI"
# url = "https://thevasudevgupta.github.io/"


def get_doc(file_path: str, cache_dir=None):
    if cache_dir is None:
        cache_dir = "~/.cache/tvgbot"
    cache_dir = Path(cache_dir).expanduser()
    cache_dir.mkdir(exist_ok=True, parents=True)

    file_type = get_file_type(file_path)

    if is_http_url(file_path):
        # if file_path is a url, we download it to cache dir
        response = make_request(file_path)
        file_path = cache_dir / hash_sha256(file_path)
        with open(file_path, "wb") as f:
            f.write(response.content)

    if file_type == "pdf":
        doc = parse_pdf(file_path)
    elif file_type == "docx":
        doc = parse_word(file_path)
    elif file_type == "pptx":
        doc = parse_ppt(file_path)
    elif file_type == "txt":
        doc = parse_txt(file_path)
    elif file_type == "html":
        doc = parse_html_bs(file_path)
    elif file_type == "csv":
        doc = parse_csv(file_path)
    elif file_type == "tsv":
        doc = parse_tsv(file_path)
    elif file_type in {"xlsx", "xls"}:
        doc = parse_excel(file_path)
    else:
        raise ValueError(f"File Type: {file_type} is not supported.")

    return get_plain_doc(doc)
